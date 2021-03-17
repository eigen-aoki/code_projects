# This script automates the Planning Document Creation Process
# This script is updated with instead of asking the user to input, we can automate using a form

import os
from os import path
from openpyxl import load_workbook
import docx
from mailmerge import MailMerge
import pptx
from pptx import Presentation
from pptx.util import Pt
from dateutil.parser import parse
from datetime import timedelta
import pandas as pd
import datetime
import shutil
from tqdm import tqdm
import time


print("Can you please tell me where the form and base templates are located?")
print("Just copy and paste the full address bar")
print("e.g. C:\\Documents\\Form")

form_location = input("Form Location: ")

os.chdir(form_location)


print("\nCan you please tell me the file name?")
print("e.g. Form.xlsx")
form_name = input("File Name: ")

time.sleep(1)
print("\nOK - Thanks. Let's Get Started")

time.sleep(1)

panda_df = pd.read_excel(form_name, sheet_name = "Input")

column_name_list = list(panda_df.columns)


clean_name_list = column_name_list[3:]

#Define the planning maker function
def planning_maker(input_value_list):

    licensee_full_name = input_value_list[0]
    licensee_short_name = input_value_list[1]

    #Turn the Licensee Contact List into String for the PPT
    licensee_Contacts= input_value_list[2]
    licensee_Contacts_split = licensee_Contacts.split(";")
    licensee_Contacts_string = str('\n'.join(licensee_Contacts_split))


    #Update Audit Start & End period to allow for different formattings
    audit_start_period = input_value_list[3]

    audit_end_period = input_value_list[4]
    print(audit_end_period)


    notification_date = input_value_list[5]
    ko_call_date = input_value_list[6]

    #Turn the Audit Team List into String for the PPT
    audit_team = input_value_list[7]
    audit_team_split = audit_team.split(";")
    audit_team_string = str('\n'.join(audit_team_split))

    proposed_fieldwork_date = input_value_list[8]

    licensee_signature = licensee_full_name

    #Set the save file names based on the short licensee name
    nda_doc_name = "CCC_" + licensee_short_name + "_NDA.docx"
    ko_deck_name = "Client_" + licensee_short_name + "_KO Presentation.pptx"
    drl_file_name = "Client_" + licensee_short_name + "_DRL.xlsx"

    #Set Dates for Phases of Audit
    #Set function for Getting the next Monday
    def get_next_monday(year, month, day):
        date0 = datetime.date(year, month, day)
        next_monday = date0 + datetime.timedelta(7 - date0.weekday() or 7)
        return next_monday

    strng_phase2_start_date = ""
    if isinstance(ko_call_date, datetime.date) == True:
        phase2_start_date = get_next_monday(ko_call_date.year, ko_call_date.month, ko_call_date.day)
        strng_phase2_start_date = phase2_start_date.strftime("%b %d")
    else: strng_phase2_start_date = "TBD"

    strng_initiation_end_date = ""
    if isinstance(ko_call_date, datetime.date) == True:
        strng_initiation_end_date = ko_call_date.strftime("%b %d")
    else: strng_initiation_end_date = "TBD"

    #Fieldwork Start date is first monday after phase 2 end dates
    fieldwork_start_date = proposed_fieldwork_date
    #fieldwork end date is 4 days after fieldwork start dates
    fieldwork_end_date = fieldwork_start_date + timedelta(days = 4)

    #data collection phase is basedon fieldwork Dates going backwards
    phase2_end_date = fieldwork_start_date - timedelta(days = 14)

    #Wrap up start date is first mondya after Fieldwork
    wrapup_start_date = get_next_monday(fieldwork_end_date.year, fieldwork_end_date.month, fieldwork_end_date.day)
    #wrap up end date is 3 weeks 19 days after wrap up
    wrapup_end_date = wrapup_start_date + timedelta(days = 25)

    #Gather string versions of the dates
    strng_initiation_start_date = notification_date.strftime("%b %d")
    strng_phase2_end_date = phase2_end_date.strftime("%b %d")
    strng_phase2_end_date_full = phase2_end_date.strftime("%B %d")
    strng_fieldwork_start_date = fieldwork_start_date.strftime("%b %d")
    strng_fieldwork_end_date = fieldwork_end_date.strftime("%b %d")
    strng_wrapup_start_date = wrapup_start_date.strftime("%b %d")
    strng_wrapup_end_date = wrapup_end_date.strftime("%b %d")

    output_location = form_location + "\\" + licensee_short_name
    #If the output location does not exist, create directory
    if os.path.exists(output_location):
        pass
    else:
        os.makedirs(licensee_short_name)

    #Create a function to add the correct values in the NDA Template
    #uses mailmerge to populate merge fields
    # print("Creating the NDA. Give me a moment........\n")
    nda_template = 'NDA_Template.docx'
    with MailMerge(nda_template) as document:
        document.merge(Licensee_Name_Full = licensee_full_name, Licensee_Full_Signature = licensee_signature)
        document.write(nda_doc_name)
    # print(f"\nFinished the NDA. It's saved as {nda_doc_name}. \nLet's do the DRL next \n")
    shutil.move(nda_doc_name, output_location)

    #Create a formula which adds the correct values in the DRL Template
    # print("Creating the DRL. Give me a moment........")
    drl_template_workbook = load_workbook('DRL_Template.xlsx')
    drl_template_worksheet = drl_template_workbook.active


    drl_template_worksheet['A2'] = licensee_full_name + " (\"Licensee\")"
    drl_template_worksheet['A4'] = "Audit Period: " + audit_start_period + " - " + audit_end_period

    drl_template_worksheet.sheet_view.showGridLines
    drl_template_worksheet.sheet_view.showGridLines = False
    drl_template_worksheet.sheet_view.zoomScale = 80

    drl_template_workbook.save(drl_file_name)
    # print(f"\nFinished the DRL. It's saved as \"{drl_file_name}\"\n")
    shutil.move(drl_file_name, output_location)

    # print("Finally, let's do the KO Deck\n")

    prs = Presentation('KO_Template.pptx')
    slides = [slide for slide in prs.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)

    #Create a formula which adds the correct values in the KO Deck
    def replace_text(presentation, replacements: dict, shapes: list):
        slides = [slide for slide in prs.slides]
        shapes = []
        for slide in slides:
            for shape in slide.shapes:
                shapes.append(shape)
        for shape in shapes:
            for match, replacement in replacements.items():
                if shape.has_text_frame:
                    if (shape.text.find(match)) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                cur_text = run.text
                                new_text = cur_text.replace(str(match), str(replacement))
                                run.text = new_text

    replace_text(prs, {'Insert_Licensee_Full': licensee_full_name}, shapes)
    replace_text(prs, {'Insert_Licensee_Short': licensee_short_name}, shapes)
    replace_text(prs, {'Insert_Licensee_Contact': licensee_Contacts_string}, shapes)
    replace_text(prs, {'Insert_Audit_Team': audit_team_string}, shapes)
    replace_text(prs, {'Insert_Audit_Start': audit_start_period}, shapes)
    replace_text(prs, {'Insert_Audit_End': audit_end_period}, shapes)
    replace_text(prs, {'Insert_Init_Date': strng_initiation_start_date}, shapes)
    replace_text(prs, {'Insert_Init_End': strng_initiation_end_date}, shapes)
    replace_text(prs, {'Insert_ph2_start': strng_phase2_start_date}, shapes)
    replace_text(prs, {'Insert_ph2_end': strng_phase2_end_date}, shapes)
    replace_text(prs, {'Insert_fldwk_strt': strng_fieldwork_start_date}, shapes)
    replace_text(prs, {'Insert_fldwk_end': strng_fieldwork_end_date}, shapes)
    replace_text(prs, {'Insert_wrp_strt': strng_wrapup_start_date}, shapes)
    replace_text(prs, {'Insert_wrp_end': strng_wrapup_end_date}, shapes)

    # print(f"Finished the KO Deck. It's saved as \"{ko_deck_name}\" \n")
    prs.save(ko_deck_name)
    shutil.move(ko_deck_name, output_location)

#iterate over all of the clean_names
for i in tqdm(range(len(clean_name_list))):
    try:
        licensee = clean_name_list[i]
        licensee_value_list = panda_df[licensee].tolist()
        planning_maker(licensee_value_list)
        time.sleep(0.1)
    except:
        pass

#Print Finished Message
print("All done - Bye Bye")
